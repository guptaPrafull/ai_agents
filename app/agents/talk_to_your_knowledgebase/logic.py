import openai
from pypdf import PdfReader
import docx
from pptx import Presentation
from app.config import settings
import os
import uuid
import tempfile
from fastapi import UploadFile

openai.api_key = settings.OPENAI_API_KEY

async def save_upload_file(file: UploadFile) -> str:
    temp_dir = tempfile.gettempdir()

    original_filename = file.filename
    extension = os.path.splitext(original_filename)[-1].lower()

    print(f"Received file: {original_filename}")
    print(f"Extracted extension: {extension}")

    if extension not in [".pdf", ".docx", ".pptx"]:
        raise ValueError("Unsupported file type. Use .pdf, .docx, or .pptx")

    filename = f"{uuid.uuid4().hex}{extension}"
    file_path = os.path.join(temp_dir, filename)

    with open(file_path, "wb") as buffer:
        content = await file.read()
        buffer.write(content)

    print(f"Saved file path: {file_path}")
    return file_path

# ---------------------------------------------
# Extract Text from PDF
# ---------------------------------------------
def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            try:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            except Exception as e:
                print(f"Could not extract text from a page: {e}")
        return text.strip()
    except Exception as e:
        raise ValueError(f"Failed to read PDF: {e}")


# ---------------------------------------------
# Extract Text from Any Supported Document Type
# ---------------------------------------------
def extract_text_from_file(file_path):
    if file_path.lower().endswith(".pdf"):
        return extract_text_from_pdf(file_path)

    elif file_path.lower().endswith(".docx"):
        try:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            raise ValueError(f"Failed to read DOCX: {e}")

    elif file_path.lower().endswith(".pptx"):
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise ValueError(f"Failed to read PPTX: {e}")

    else:
        raise ValueError("Unsupported file type. Use .pdf, .docx, or .pptx")


# ---------------------------------------------
# Summarize Document
# ---------------------------------------------
def summarize_document(document_text):
    prompt = f"""
You are a professional summarizer for technical presentations and educational slide decks.

Carefully analyze the document below and generate a **detailed, structured summary** that includes:

1. A clear, engaging 3–4 sentence **overview** describing the topic, objective, and target audience of the document.
2. All **major sections and subtopics**, presented with meaningful **emoji-enhanced headings**.
3. **Bullet points** under each section that preserve specific technical terms, types, techniques, examples, and challenges.
4. A final **Conclusion** section with key takeaways or forward-looking remarks.

Your output must:
- Fully reflect the **entire content** of the document.
- Be professional, precise, and scannable using Markdown-style formatting.
- NOT add speculative content or mention "expected details."
- Use the original structure and topics from the document as faithfully as possible.

---

Document:
\"\"\"
{document_text}
\"\"\"

---

###  **Overview**
(A 3–4 sentence summary paragraph.)

---

###  **[Section Title]**
- Bullet
- Bullet

...

###  **Conclusion**
- Bullet summarizing key takeaways
- Bullet on relevance, benefits, or future outlook
"""

    response = openai.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You summarize documents clearly and professionally."},
            {"role": "user", "content": prompt.strip()}
        ],
        temperature=0.3,
        max_tokens=500
    )
    return response.choices[0].message.content.strip()


# ---------------------------------------------
# Generate FAQ by Topic
# ---------------------------------------------
def generate_faq_by_topic(topic, content):
    prompt = f"""
You are a professional FAQ writer.

Your task is to generate a list of 3–5 Frequently Asked Questions (FAQs) specifically focused on the following topic:

📄 **Topic**: {topic}

Instructions:
- Extract or infer all questions and answers strictly from the provided content.
- Do not invent any information or go beyond the scope of the content.
- Only generate questions that are directly related to the topic.
- Ensure all answers are concise, accurate, and based solely on the content.
- If the content does not provide enough information, generate only the questions and answers that can be clearly supported.
- Present the FAQ in a clean, professional, and easy-to-read format.

---

Content:
\"\"\"
{content}
\"\"\"
"""

    response = openai.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You generate professional, topic-specific FAQs from documents."},
            {"role": "user", "content": prompt.strip()}
        ],
        temperature=0.3,
        max_tokens=900
    )
    return response.choices[0].message.content.strip()


# ---------------------------------------------
# Chat With Document - Returns Updated History
# ---------------------------------------------
def chat_with_document(document_text, user_input, history):
    if not history:
        history = [
            {"role": "system",
             "content": "You are a helpful assistant that answers questions based only on the provided document."},
            {"role": "user", "content": f"Document:\n\"\"\"\n{document_text}\n\"\"\""}
        ]

    history.append({"role": "user", "content": user_input})

    response = openai.chat.completions.create(
        model="gpt-4",
        messages=history,
        temperature=0.3,
        max_tokens=700
    )

    reply = response.choices[0].message.content.strip()
    history.append({"role": "assistant", "content": reply})

    return reply, history
